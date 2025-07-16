from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns

def create_hist_plot(df, column, title, xlabel, ylabel, color='#2ABB00'):
    sns.set_style('whitegrid')
    plt.figure(figsize=(6,4))
    ax = sns.histplot(df[column], bins=20, kde=True, color=color)
    ax.set_title(title, fontsize=14, weight='bold', color=color)
    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=150)
    plt.close()
    buf.seek(0)
    return buf

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_text_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.text = content

def add_picture_slide(prs, title, image_stream):
    slide_layout = prs.slide_layouts[5]  # Title and Content (Picture)
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture(image_stream, left, top, width=Inches(8))

def generate_presentation(csv_path, pptx_path):
    df = pd.read_csv(csv_path)
    prs = Presentation()

    # Заглавный слайд
    add_title_slide(prs, "Отчёт по вкладам Сбербанка", "Автоматически сгенерированная презентация")

    # Слайд со статистикой
    content = f"Общее количество записей: {len(df)}"
    add_text_slide(prs, "Общая статистика", content)

    # Слайд с распределением ставок
    rate_plot = create_hist_plot(df, 'rate', 'Распределение процентных ставок', 'Ставка (%)', 'Количество вкладов')
    add_picture_slide(prs, "Распределение ставок", rate_plot)

    # Слайд с распределением сроков
    term_plot = create_hist_plot(df, 'term_months', 'Распределение сроков (месяцы)', 'Срок (мес.)', 'Количество вкладов')
    add_picture_slide(prs, "Распределение сроков", term_plot)

    prs.save(pptx_path)
    print(f"Презентация сохранена в {pptx_path}")

if __name__ == "__main__":
    generate_presentation("../data/clean/clean_deposits.csv", "presentation.pptx")